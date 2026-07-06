import streamlit as st
from pptx import Presentation
import pandas as pd
import numpy as np
import warnings
from sklearn.ensemble import RandomForestClassifier

# 경고 메시지 끄기
warnings.filterwarnings('ignore', category=UserWarning)

# 웹페이지 상단 제목 및 레이아웃 설정
st.set_page_config(page_title="한의학 체질 기반 치료 예측 AI", layout="centered")
st.image("logo.png", use_container_width=True)
st.title("사상체질 기반 AI 치료 효과 예측 시스템")
st.markdown("환자의 사상체질과 라이프스타일 데이터를 기반으로 맞춤형 치료 결과를 예측하는 인공지능 프로토타입입니다.")

# ==========================================
# 1. 실제 PPTX 데이터 로드 및 AI 모델 학습 (캐싱 처리)
# ==========================================
@st.cache_resource
def train_ai_model():
    pptx_path = "환자데이터.pptx" 
    all_rows = []
    columns = ['sex', 'age', 'constitution', 'sleep', 'stress', 'exercise', 'smoking', 'alcohol', 'treatment']
    
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        clean_row = [cell.text_frame.text.strip() for cell in row.cells]
                        
                        # 제목 줄 건너뛰기
                        if row_idx == 0 and any(keyword in str(clean_row[0]) or keyword in str(clean_row[1]) for keyword in ['sex', '성별', 'ID', '번호', 'no']):
                            continue
                        
                        # 첫 번째 칸(환자 넘버) 제외한 9개 데이터 추출
                        actual_data = clean_row[1:]
                        
                        # 숫자 정제
                        fixed_row = []
                        for cell_text in actual_data:
                            cleaned_num = "".join([c for c in cell_text if c.isdigit() or c == '.'])
                            fixed_row.append(cleaned_num if cleaned_num != "" else "0")
                        
                        if len(fixed_row) < 9:
                            while len(fixed_row) < 9: fixed_row.append("1")
                        else:
                            fixed_row = fixed_row[:9]
                            
                        all_rows.append(fixed_row)
        
        # 데이터프레임 구축 및 숫자 변환
        df = pd.DataFrame(all_rows, columns=columns)
        df = df.apply(pd.to_numeric, errors='coerce').dropna()
        
        # AI 모델 학습
        X = df[columns[:-1]]
        y = df['treatment'].astype(int)
        model = RandomForestClassifier(n_estimators=100, random_state=42)
        model.fit(X, y)
        
        return model, len(df), "success"
        
    except Exception as e:
        np.random.seed(42)
        num_samples = 1000
        mock_df = pd.DataFrame({
            'sex': np.random.choice([0, 1], num_samples),
            'age': np.random.randint(20, 80, num_samples),
            'constitution': np.random.choice([0, 1, 2, 3], num_samples),
            'sleep': np.random.uniform(5, 9, num_samples),
            'stress': np.random.randint(1, 11, num_samples),
            'exercise': np.random.uniform(0, 3, num_samples),
            'smoking': np.random.choice([0, 1], num_samples),
            'alcohol': np.random.choice([0, 1], num_samples),
            'treatment': np.random.choice([0, 1, 2], num_samples)
        })
        X = mock_df[columns[:-1]]
        y = mock_df['treatment'].astype(int)
        model = RandomForestClassifier(n_estimators=100, random_state=42)
        model.fit(X, y)
        return model, num_samples, f"fallback ({e})"

model, data_count, status = train_ai_model()

if "success" in status:
    st.success(f"📊 실제 PPTX 파일에서 {data_count}명의 환자 데이터를 성공적으로 읽어와 AI 학습을 완료했습니다!")
else:
    st.warning(f"⚠️ 실제 PPTX 파일을 읽지 못해 가상 데이터({data_count}명)로 시뮬레이션 중입니다. (오류: {status})")

# ==========================================
# 2. 웹 UI 화면 구성 (좌측 사이드바 입력창)
# ==========================================
st.sidebar.header("📋 신규 환자 정보 입력")

sex = st.sidebar.selectbox("성별", options=[0, 1], format_func=lambda x: "남성 (0)" if x == 0 else "여성 (1)")
age = st.sidebar.slider("나이 (세)", min_value=0, max_value=100, value=45)
constitution = st.sidebar.selectbox("사상 체질", options=[0, 1, 2, 3], format_func=lambda x: ["태양인 (0)", "태음인 (1)", "소양인 (2)", "소음인 (3)"][x])
sleep = st.sidebar.slider("하루 수면 시간 (시간)", min_value=0.0, max_value=24.0, value=7.0, step=0.5)
stress = st.sidebar.slider("스트레스 지수 (1~10)", min_value=1, max_value=10, value=5)
exercise = st.sidebar.slider("주간 운동 시간 (시간)", min_value=0.0, max_value=24.0, value=2.0, step=0.5)

# ✨ 오타 완벽 제거: 흡연 여부 / 음주 여부
smoking = st.sidebar.selectbox("흡연 여부", options=[0, 1], format_func=lambda x: "비흡연 (0)" if x == 0 else "흡연 (1)")
alcohol = st.sidebar.selectbox("음주 여부", options=[0, 1], format_func=lambda x: "비음주 (0)" if x == 0 else "음주 (1)")

# ==========================================
# 3. 실시간 예측결과 및 그래프 시각화 출력 (우측 본문)
# ==========================================
st.subheader("📊 실시간 분석 결과")

const_names = ["태양인", "태음인", "소양인", "소음인"]
summary_df = pd.DataFrame({
    "지표": ["성별", "나이", "사상체질", "수면시간", "스트레스", "운동시간", "흡연", "음주"],
    "입력값": ["여성" if sex==1 else "남성", f"{age}세", const_names[constitution], f"{sleep}시간", f"{stress}점", f"{exercise}시간", "흡연" if smoking==1 else "비흡연", "음주" if alcohol==1 else "비음주"]
})
st.dataframe(summary_df, use_container_width=True, hide_index=True)

if st.button("🔮 치료 결과 실시간 예측하기", type="primary", use_container_width=True):
    new_input = pd.DataFrame([[sex, age, constitution, sleep, stress, exercise, smoking, alcohol]], 
                             columns=['sex', 'age', 'constitution', 'sleep', 'stress', 'exercise', 'smoking', 'alcohol'])
    
    pred_code = model.predict(new_input)[0]
    pred_proba = model.predict_proba(new_input)[0]
    
    # ✨ 오타 완벽 제거: 악화 / 변화없음 / 호전됨
    result_map = {0: '🔴 악화될 가능성이 높음', 1: '🟡 변화가 없을 가능성이 높음', 2: '🟢 호전될 가능성이 높음'}
    
    st.markdown("---")
    st.markdown(f"### 🎯 AI 최종 예상 결과")
    st.markdown(f"#### 본 환자는 치료 후 **{result_map[pred_code]}**로 예측됩니다.")
    
    st.markdown("### 📈 결과별 예측 확률 분포")
    proba_df = pd.DataFrame({
        '치료 결과': ['악화', '변화없음', '호전됨'],
        '확률 (%)': [pred_proba[0]*100, pred_proba[1]*100, pred_proba[2]*100]
    })
    st.bar_chart(data=proba_df, x='치료 결과', y='확률 (%)', color="#4A90E2")
